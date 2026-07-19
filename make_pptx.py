"""
Generates a 15-slide PowerPoint presentation for the HR Management System project.
Run: python make_pptx.py
Output: HR_Management_System.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Brand colours ────────────────────────────────────────────────────────────
C_PRIMARY   = RGBColor(0x4F, 0x46, 0xE5)   # Indigo
C_ACCENT    = RGBColor(0x06, 0xB6, 0xD4)   # Cyan
C_DARK      = RGBColor(0x0F, 0x17, 0x2A)   # Dark navy
C_TEXT      = RGBColor(0x1E, 0x29, 0x3B)   # Slate 800
C_SUBTEXT   = RGBColor(0x64, 0x74, 0x8B)   # Slate 500
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xF8, 0xFA, 0xFF)
C_EMERALD   = RGBColor(0x10, 0xB9, 0x81)
C_AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
C_ROSE      = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.name  = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=16, color=None, bold_first=False,
                   line_spacing=1.2, font_name="Calibri"):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet character
        run = p.add_run()
        run.text = ("▶  " if idx == 0 and bold_first else "•  ") + bullet
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.bold = (idx == 0 and bold_first)
        if color:
            run.font.color.rgb = color
    return txBox


def card(slide, left, top, width, height, fill=None, accent_left=None, radius=False):
    """White card with optional left accent bar."""
    bg_color = fill or C_WHITE
    rect = add_rect(slide, left, top, width, height, fill_rgb=bg_color, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    if accent_left:
        add_rect(slide, left, top, 0.04, height, fill_rgb=accent_left)
    return rect


def section_header(slide, text, color=C_PRIMARY):
    add_rect(slide, 0, 0, 13.33, 0.06, fill_rgb=color)
    add_rect(slide, 0, 7.44, 13.33, 0.06, fill_rgb=color)
    add_text(slide, text.upper(), 0.4, 0.12, 12, 0.35,
             font_size=8, bold=True, color=color, font_name="Calibri")


def slide_bg(slide, dark=False):
    color = C_DARK if dark else C_LIGHT_BG
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=color)


# ─── SLIDE 1 – Title ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, dark=True)

# Gradient bar (simulate with two rects)
add_rect(s, 0, 0, 6.67, 7.5, fill_rgb=C_PRIMARY)
add_rect(s, 6.67, 0, 6.66, 7.5, fill_rgb=RGBColor(0x06, 0x3D, 0x5E))

# Decorative circles
add_rect(s, -1, -1, 5, 5, fill_rgb=RGBColor(0x5E, 0x52, 0xF5))          # big circle (clipped)
add_rect(s, 10, 5, 4, 4, fill_rgb=RGBColor(0x06, 0x8B, 0xA4))

# Product icon box
add_rect(s, 1.2, 1.6, 0.7, 0.7, fill_rgb=RGBColor(0x38, 0x2E, 0xF5))
add_text(s, "HR", 1.22, 1.62, 0.66, 0.66, font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(s, "AI-HRMS", 2.1, 1.6, 5, 0.7, font_size=28, bold=True, color=C_WHITE)
add_text(s, "AI-Powered Human Resource", 1.2, 2.5, 9, 0.6, font_size=32, bold=True, color=C_WHITE)
add_text(s, "Management System", 1.2, 3.15, 9, 0.6, font_size=32, bold=True, color=C_ACCENT)
add_text(s, "React · TypeScript · Flask · Machine Learning", 1.2, 3.9, 10, 0.45,
         font_size=15, color=RGBColor(0xBF, 0xDB, 0xFF))
add_text(s, "Full-stack HR platform with predictive analytics and a responsive dashboard",
         1.2, 4.5, 9, 0.6, font_size=13, italic=True,
         color=RGBColor(0x93, 0xC5, 0xFD))

# Right side info tags
for i, tag in enumerate(["20 Employees", "7 Departments", "ML Predictions", "REST API"]):
    add_rect(s, 8.5, 2.3 + i * 0.8, 3.5, 0.55, fill_rgb=RGBColor(0xFF, 0xFF, 0xFF, ))
    add_text(s, "✦  " + tag, 8.65, 2.35 + i * 0.8, 3.2, 0.45, font_size=13,
             bold=True, color=C_PRIMARY)

add_text(s, "2026", 1.2, 6.8, 4, 0.4, font_size=11,
         color=RGBColor(0x93, 0xC5, 0xFD))

# ─── SLIDE 2 – Agenda / Table of Contents ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Agenda")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Table of Contents", 0.5, 0.6, 12, 0.85,
         font_size=30, bold=True, color=C_WHITE)

items = [
    ("01", "Project Overview",          "Purpose, objectives, and high-level scope"),
    ("02", "Architecture & Stack",      "Frontend, backend, data flow"),
    ("03", "Frontend Features",         "UI modules and page breakdown"),
    ("04", "Backend & API",             "Flask REST endpoints and data layer"),
    ("05", "AI Prediction Engine",      "ML model, inputs, outputs"),
    ("06", "Authentication",            "Role-based auth and session strategy"),
    ("07", "Employee Management",       "List, detail, onboarding flows"),
    ("08", "Attendance & Performance",  "Tracking and review dashboards"),
    ("09", "Analytics & Reporting",     "Charts, KPIs, exports"),
    ("10", "Notifications & Settings",  "Notification center, theme, preferences"),
    ("11", "Deployment",                "Vercel setup and serverless config"),
    ("12", "Data & Persistence",        "JSON datastore and seed data"),
    ("13", "Mock vs API Behavior",      "Current data source per module"),
    ("14", "Known Limitations",         "Auth, persistence, coverage gaps"),
    ("15", "Roadmap",                   "Upcoming improvements"),
]

cols = [items[:8], items[8:]]
col_x = [0.4, 6.9]
for ci, col in enumerate(cols):
    for ri, (num, title, desc) in enumerate(col):
        y = 1.75 + ri * 0.68
        add_rect(s, col_x[ci], y, 0.45, 0.45, fill_rgb=C_PRIMARY)
        add_text(s, num, col_x[ci], y + 0.04, 0.45, 0.4, font_size=11, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, col_x[ci] + 0.55, y, 2.6, 0.3,
                 font_size=12, bold=True, color=C_TEXT)
        add_text(s, desc, col_x[ci] + 0.55, y + 0.26, 2.6, 0.32,
                 font_size=9, italic=True, color=C_SUBTEXT)

# ─── SLIDE 3 – Project Overview ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Project Overview")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Project Overview", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 03 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

add_text(s, "What is AI-HRMS?", 0.5, 1.75, 12, 0.45, font_size=18, bold=True, color=C_PRIMARY)
add_text(s,
    "AI-HRMS is a full-stack Human Resource Management System that centralises workforce operations "
    "into a single, modern dashboard. It combines a React/TypeScript frontend with a Python Flask "
    "backend and a machine learning prediction engine to surface actionable HR insights.",
    0.5, 2.25, 12.3, 0.9, font_size=13, color=C_TEXT)

# Three objective cards
obj_cards = [
    (C_PRIMARY, "Centralise Operations", "One platform for employees, attendance, performance, and reporting."),
    (C_EMERALD, "Predictive Intelligence", "Random Forest ML model classifies employee performance proactively."),
    (C_ACCENT,  "Demo-Ready Resilience",  "Graceful mock-data fallback keeps the UI functional without a backend."),
]
for i, (col, title, body) in enumerate(obj_cards):
    x = 0.4 + i * 4.25
    add_rect(s, x, 3.3, 4.0, 2.4, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_rect(s, x, 3.3, 4.0, 0.08, fill_rgb=col)
    add_text(s, title, x + 0.18, 3.5, 3.7, 0.45, font_size=14, bold=True, color=col)
    add_text(s, body, x + 0.18, 4.0, 3.65, 1.5, font_size=12, color=C_TEXT)

add_text(s, "Target users: HR Managers · Administrators · Employees", 0.5, 5.95, 12, 0.35,
         font_size=11, italic=True, color=C_SUBTEXT)

# ─── SLIDE 4 – Architecture & Tech Stack ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Architecture & Stack")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Architecture & Tech Stack", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 04 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# Diagram boxes
layers = [
    ("BROWSER / CLIENT", C_PRIMARY,
     "React 19 · TypeScript · Vite\nTailwind CSS · Framer Motion · Recharts\nAxios · React Hook Form · React Hot Toast"),
    ("REST API LAYER", C_ACCENT,
     "Flask · Flask-CORS\n/api/* routes · JSON responses\nCORS enabled for all origins"),
    ("DATA + ML LAYER", C_EMERALD,
     "db.json (JSON datastore)\nHRDatabase class · Auto-seeded\nRandom Forest Classifier · trained_model.pkl"),
]
for i, (label, col, detail) in enumerate(layers):
    y = 1.7 + i * 1.7
    add_rect(s, 1.0, y, 11.33, 1.5, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_rect(s, 1.0, y, 0.08, 1.5, fill_rgb=col)
    add_rect(s, 1.08, y, 2.6, 1.5, fill_rgb=col)
    add_text(s, label, 1.15, y + 0.5, 2.45, 0.6, font_size=11, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, detail, 3.8, y + 0.2, 8.3, 1.1, font_size=12, color=C_TEXT)
    # Arrow down
    if i < 2:
        add_text(s, "▼", 6.4, y + 1.5, 0.7, 0.4, font_size=14, bold=True,
                 color=C_PRIMARY, align=PP_ALIGN.CENTER)

add_text(s, "Frontend communicates with /api/* over HTTP · Bearer token from localStorage user.id",
         0.5, 6.95, 12.3, 0.35, font_size=10, italic=True, color=C_SUBTEXT)

# ─── SLIDE 5 – Frontend Features & Pages ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Frontend Features")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Frontend Features & Pages", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 05 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

pages = [
    ("🔐", "Auth",           "Login · Register · Forgot Password\nRole selector: Admin / HR / Employee"),
    ("📊", "Dashboard",      "KPI stat cards · Area & radar charts\nActivity stream · Quick action buttons"),
    ("👥", "Employees",      "Searchable list · Detail profile\nAdd employee form · Delete action"),
    ("🏢", "Departments",    "7 department cards · Budget progress\nHeadcount · Productivity scores"),
    ("📅", "Attendance",     "Status tracker · Monthly trend\nLine & bar charts · Log table"),
    ("⭐", "Performance",    "KPI radar · Department scores\nReview table · Recommendation badge"),
    ("🤖", "AI Prediction",  "Input sliders · ML call / fallback\nConfidence ring · Feature importance"),
    ("📈", "Analytics",      "Growth · Gender · Pie charts\nSalary · Heatmap placeholder"),
    ("📄", "Reports",        "Overview / Turnover / Performance\nDepartment · Export PDF/CSV"),
    ("🔔", "Notifications",  "Filter by type & read state\nMark read · Delete · Count badges"),
]
cols_per_row = 5
for i, (icon, name, detail) in enumerate(pages):
    col = i % cols_per_row
    row = i // cols_per_row
    x = 0.3 + col * 2.54
    y = 1.75 + row * 2.4
    add_rect(s, x, y, 2.35, 2.15, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_text(s, icon, x + 0.1, y + 0.12, 0.5, 0.5, font_size=18, align=PP_ALIGN.LEFT)
    add_text(s, name, x + 0.65, y + 0.15, 1.6, 0.4, font_size=12, bold=True, color=C_PRIMARY)
    add_text(s, detail, x + 0.1, y + 0.6, 2.15, 1.45, font_size=9, color=C_TEXT)

# ─── SLIDE 6 – Backend & API Layer ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Backend & API")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Backend & REST API", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 06 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

add_text(s, "Two Backend Variants", 0.5, 1.72, 5, 0.4, font_size=16, bold=True, color=C_PRIMARY)

variants = [
    ("backend/app.py", C_PRIMARY,
     "Port 5678\nIncludes PDF + CSV export endpoints\nIncludes train_model.py script\nRecommended for local development"),
    ("api/index.py", C_ACCENT,
     "Port 5000\nVercel serverless-compatible imports\nIdentical core API surface\nDeployed as Python runtime on Vercel"),
]
for i, (name, col, desc) in enumerate(variants):
    x = 0.4 + i * 6.2
    add_rect(s, x, 2.15, 5.85, 2.1, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_rect(s, x, 2.15, 5.85, 0.08, fill_rgb=col)
    add_text(s, name, x + 0.18, 2.28, 5.5, 0.4, font_size=13, bold=True, color=col)
    add_text(s, desc, x + 0.18, 2.72, 5.5, 1.4, font_size=12, color=C_TEXT)

add_text(s, "Endpoint Groups", 0.5, 4.45, 5, 0.4, font_size=16, bold=True, color=C_PRIMARY)

groups = [
    ("/auth",         "login · register · logout · me"),
    ("/employees",    "CRUD with id, avatar, status, salary"),
    ("/departments",  "CRUD + auto headcount sync"),
    ("/attendance",   "check-in · check-out · by date/employee"),
    ("/performance",  "CRUD · per-employee · score calc"),
    ("/predictions",  "predict · predict/predict · per-employee"),
    ("/analytics",    "dashboard · growth · trend · distribution"),
    ("/notifications","list · read · read-all · delete"),
    ("/reports",      "employees · attendance · performance · predictions · export/pdf · export/excel"),
]
col_data = [groups[:5], groups[5:]]
for ci, col_list in enumerate(col_data):
    for ri, (ep, detail) in enumerate(col_list):
        x = 0.4 + ci * 6.5
        y = 4.9 + ri * 0.42
        add_rect(s, x, y, 1.5, 0.35, fill_rgb=C_PRIMARY)
        add_text(s, ep, x + 0.05, y + 0.03, 1.4, 0.3, font_size=9, bold=True, color=C_WHITE)
        add_text(s, detail, x + 1.6, y + 0.03, 4.6, 0.3, font_size=9, color=C_TEXT)

# ─── SLIDE 7 – AI Prediction Engine ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "AI Prediction Engine")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "AI Prediction Engine", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 07 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# Left: algorithm details
add_rect(s, 0.4, 1.72, 5.8, 5.0, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 0.4, 1.72, 5.8, 0.08, fill_rgb=C_PRIMARY)
add_text(s, "Random Forest Classifier", 0.6, 1.85, 5.4, 0.4, font_size=14, bold=True, color=C_PRIMARY)
algo_lines = [
    "Algorithm: RandomForestClassifier (scikit-learn)",
    "Estimators: 100 decision trees",
    "Max depth: 10  |  min_samples_split: 5",
    "Training data: 500 synthetic samples",
    "Weighted score threshold: ≥80 → High, ≥60 → Medium",
    "Model file: trained_model.pkl (joblib)",
    "",
    "Input features (all 0–100 scale):",
    "  • Attendance Score",
    "  • Experience Score",
    "  • KPI Score",
    "  • Training Completion Score",
    "  • Leave Impact (days)",
    "  • Previous Rating",
]
for i, line in enumerate(algo_lines):
    add_text(s, line, 0.6, 2.3 + i * 0.33, 5.4, 0.32,
             font_size=10.5, bold=line.startswith("Input"),
             color=C_TEXT if not line.startswith("  •") else C_SUBTEXT)

# Right: output classes
add_rect(s, 6.6, 1.72, 6.4, 2.1, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 6.6, 1.72, 6.4, 0.08, fill_rgb=C_EMERALD)
add_text(s, "Output Classes", 6.8, 1.85, 6.0, 0.4, font_size=14, bold=True, color=C_EMERALD)
for i, (cls, col, example) in enumerate([
    ("High Performer",   C_EMERALD, "Score ≥80  |  Conf ~94%"),
    ("Medium Performer", C_AMBER,   "Score 60–80 |  Conf ~78%"),
    ("Low Performer",    C_ROSE,    "Score <60  |  Conf ~89%"),
]):
    y = 2.3 + i * 0.5
    add_rect(s, 6.8, y, 0.2, 0.25, fill_rgb=col)
    add_text(s, cls, 7.1, y - 0.02, 3.0, 0.3, font_size=12, bold=True, color=col)
    add_text(s, example, 7.1, y + 0.26, 5.5, 0.28, font_size=10, color=C_SUBTEXT)

# Right: API flows
add_rect(s, 6.6, 4.0, 6.4, 2.72, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 6.6, 4.0, 6.4, 0.08, fill_rgb=C_ACCENT)
add_text(s, "API Response Payload", 6.8, 4.13, 6.0, 0.4, font_size=14, bold=True, color=C_ACCENT)
payload_lines = [
    "prediction       →  class label",
    "confidence       →  float 0–100",
    "confidenceScores →  per-class breakdown",
    "featureImportance→  top-5 with impact %",
    "recommendations  →  action items",
    "riskFactors      →  detected issues",
]
for i, pl in enumerate(payload_lines):
    add_text(s, pl, 6.8, 4.6 + i * 0.36, 6.0, 0.34, font_size=10.5, color=C_TEXT)

# ─── SLIDE 8 – Authentication ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Authentication")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Authentication & Session Management", 0.5, 0.6, 10, 0.85, font_size=28, bold=True, color=C_WHITE)
add_text(s, "Slide 08 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# Flow diagram
steps = [
    ("User fills\nLogin/Register form", C_PRIMARY),
    ("AuthContext calls\nAPI endpoint", C_ACCENT),
    ("On success:\nsave to localStorage", C_EMERALD),
    ("On failure:\nMock fallback", C_AMBER),
    ("Axios interceptor\nattaches Bearer token", C_PRIMARY),
]
for i, (label, col) in enumerate(steps):
    x = 0.5 + i * 2.5
    add_rect(s, x, 1.85, 2.1, 1.1, fill_rgb=col)
    add_text(s, label, x + 0.08, 1.95, 1.95, 0.9, font_size=10.5, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "→", x + 2.1, 2.22, 0.4, 0.4, font_size=18, bold=True,
                 color=C_PRIMARY, align=PP_ALIGN.CENTER)

# Roles
add_text(s, "Roles", 0.5, 3.2, 3, 0.4, font_size=16, bold=True, color=C_PRIMARY)
roles = [
    ("Administrator", "Full system access. Manage all data.", C_PRIMARY),
    ("HR Manager",    "Employee ops, performance, predictions.", C_ACCENT),
    ("Employee",      "View own profile and records.", C_EMERALD),
]
for i, (role, desc, col) in enumerate(roles):
    y = 3.65 + i * 0.8
    add_rect(s, 0.5, y, 5.5, 0.65, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_rect(s, 0.5, y, 0.06, 0.65, fill_rgb=col)
    add_text(s, role, 0.7, y + 0.08, 2.0, 0.35, font_size=12, bold=True, color=col)
    add_text(s, desc, 2.8, y + 0.1, 3.1, 0.45, font_size=10.5, color=C_TEXT)

# Demo accounts
add_text(s, "Demo Accounts", 6.8, 3.2, 6.0, 0.4, font_size=16, bold=True, color=C_PRIMARY)
accounts = [
    ("admin@hrms.ng", "Administrator", "U001 — Chukwuemeka Nwosu"),
    ("hr@hrms.ng",    "HR Manager",    "U002 — Fatima Al-Amin"),
    ("emp@hrms.ng",   "Employee",      "U003 — Adaeze Okonkwo"),
]
for i, (email, role, name) in enumerate(accounts):
    y = 3.65 + i * 0.8
    add_rect(s, 6.8, y, 6.1, 0.65, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_text(s, email, 7.0, y + 0.05, 2.6, 0.3, font_size=11, bold=True, color=C_PRIMARY)
    add_text(s, role, 7.0, y + 0.33, 2.0, 0.28, font_size=9.5, color=C_SUBTEXT)
    add_text(s, name, 9.7, y + 0.12, 3.0, 0.3, font_size=10, color=C_TEXT)

add_text(s, "Note: Any password is accepted in the current demo auth logic.", 0.5, 6.9, 12, 0.35,
         font_size=10, italic=True, color=C_SUBTEXT)

# ─── SLIDE 9 – Employee Management ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Employee Management")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Employee Management", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 09 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# 3 feature panels
panels = [
    ("Employee List", C_PRIMARY, [
        "20 seeded employees across 7 departments",
        "Search by name, email, ID or position",
        "Filter by department and employment status",
        "Paginated table (8 per page)",
        "Avatar initials, status badge, performance bar",
        "Delete employee (API + fallback)",
    ]),
    ("Employee Detail", C_ACCENT, [
        "Profile card with all contact & employment info",
        "KPI radar chart and 6-metric score bars",
        "Reviewer comments section",
        "API-first: employeeService.getById()",
        "Fallback to mockEmployees on API failure",
        "Performance data from latest review",
    ]),
    ("Add Employee Form", C_EMERALD, [
        "React Hook Form with validation",
        "Required: name, email, phone, position, dept",
        "Optional: salary, experience, manager, location",
        "Calls employeeService.create() on submit",
        "Falls back gracefully on API error",
        "Redirects to list on success",
    ]),
]
for i, (title, col, bullets) in enumerate(panels):
    x = 0.4 + i * 4.25
    add_rect(s, x, 1.75, 4.0, 5.0, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_rect(s, x, 1.75, 4.0, 0.08, fill_rgb=col)
    add_text(s, title, x + 0.18, 1.88, 3.7, 0.4, font_size=14, bold=True, color=col)
    for j, b in enumerate(bullets):
        add_text(s, "• " + b, x + 0.18, 2.4 + j * 0.63, 3.7, 0.56, font_size=10.5, color=C_TEXT)

# ─── SLIDE 10 – Attendance & Performance ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Attendance & Performance")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Attendance & Performance Modules", 0.5, 0.6, 10, 0.85, font_size=28, bold=True, color=C_WHITE)
add_text(s, "Slide 10 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# Attendance
add_rect(s, 0.4, 1.75, 6.0, 5.0, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 0.4, 1.75, 6.0, 0.08, fill_rgb=C_AMBER)
add_text(s, "📅  Attendance Tracker", 0.6, 1.88, 5.6, 0.4, font_size=14, bold=True, color=C_AMBER)

att_items = [
    "Stat summary: Present / Absent / Late / On Leave",
    "Monthly attendance trend (Line chart)",
    "Status breakdown (Grouped bar chart)",
    "Filterable log table with timestamps",
    "Backend: check-in / check-out endpoints",
    "Status: Present if arrival ≤09:00 else Late",
    "Hours calculated: checkOut − checkIn",
    "Currently using mock data for UI demo",
]
for i, item in enumerate(att_items):
    add_text(s, "• " + item, 0.6, 2.38 + i * 0.55, 5.7, 0.5, font_size=10.5, color=C_TEXT)

# Performance
add_rect(s, 6.9, 1.75, 6.0, 5.0, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 6.9, 1.75, 6.0, 0.08, fill_rgb=C_PRIMARY)
add_text(s, "⭐  Performance Management", 7.1, 1.88, 5.6, 0.4, font_size=14, bold=True, color=C_PRIMARY)

perf_items = [
    "10 seeded performance reviews (seed data)",
    "6 KPI dimensions per review",
    "   Productivity · Teamwork · Leadership",
    "   Communication · Innovation · Discipline",
    "Overall score = average of 6 KPIs",
    "Trend line chart (2025 monthly avg)",
    "Department horizontal bar chart",
    "Status: Completed / Pending / Overdue",
    "Recommendations: Promote / Retain / Warn / Dismiss",
    "Review score updates employee performanceScore",
]
for i, item in enumerate(perf_items):
    add_text(s, ("  " if item.startswith("   ") else "• ") + item.strip(),
             7.1, 2.38 + i * 0.44, 5.7, 0.42, font_size=10.5, color=C_TEXT)

# ─── SLIDE 11 – Analytics & Reporting ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Analytics & Reporting")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Analytics & Reporting", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 11 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# Analytics charts list
add_rect(s, 0.4, 1.72, 5.9, 5.05, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 0.4, 1.72, 5.9, 0.08, fill_rgb=C_ACCENT)
add_text(s, "📈  Analytics Page Charts", 0.6, 1.85, 5.6, 0.4, font_size=14, bold=True, color=C_ACCENT)
charts = [
    ("Employee Growth", "Area chart · 2025 monthly headcount"),
    ("Attendance Trends", "Line chart · Present / Absent / Late %"),
    ("Performance Distribution", "Pie chart · Excellent / Good / Satisfactory"),
    ("Gender Distribution", "Pie chart · Male vs Female split"),
    ("KPI Radar", "Radar chart · Company-wide metric averages"),
    ("Dept Performance vs Attendance", "Grouped bar chart"),
    ("Avg Salary by Department", "Horizontal bar · ₦ formatted"),
    ("Monthly HR Activity", "Bar chart · Hired / Resigned / Promoted"),
    ("Attendance Heatmap", "Calendar placeholder (backend integration pending)"),
]
for i, (name, desc) in enumerate(charts):
    add_text(s, name, 0.6, 2.3 + i * 0.51, 2.8, 0.3, font_size=10.5, bold=True, color=C_PRIMARY)
    add_text(s, desc, 3.5, 2.3 + i * 0.51, 2.7, 0.3, font_size=10, color=C_TEXT)

# Reports
add_rect(s, 6.7, 1.72, 6.2, 5.05, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 6.7, 1.72, 6.2, 0.08, fill_rgb=C_EMERALD)
add_text(s, "📄  Reports Page", 6.9, 1.85, 5.8, 0.4, font_size=14, bold=True, color=C_EMERALD)

rpt_items = [
    ("Report Types", "Overview · Turnover · Performance · Department"),
    ("Date Range Filter", "Week / Month / Quarter / Year"),
    ("Overview", "Employee growth line + dept bar chart"),
    ("Turnover", "Monthly hired / left / net bar chart"),
    ("Performance", "Pie distribution + summary table"),
    ("Department", "Salary and dept comparison bars"),
    ("Export: PDF", "POST /reports/export/pdf (backend/app.py)"),
    ("Export: CSV/Excel", "POST /reports/export/excel → CSV file"),
    ("UI Exports", "Buttons present; real call requires backend"),
]
for i, (label, desc) in enumerate(rpt_items):
    y = 2.3 + i * 0.51
    add_text(s, label + ":", 6.9, y, 2.1, 0.3, font_size=10.5, bold=True, color=C_PRIMARY)
    add_text(s, desc, 9.05, y, 3.7, 0.3, font_size=10, color=C_TEXT)

# ─── SLIDE 12 – Notifications & Settings ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Notifications & Settings")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Notifications & Settings", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 12 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# Notifications
add_rect(s, 0.4, 1.72, 6.0, 5.0, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 0.4, 1.72, 6.0, 0.08, fill_rgb=C_AMBER)
add_text(s, "🔔  Notification Center", 0.6, 1.85, 5.6, 0.4, font_size=14, bold=True, color=C_AMBER)
notif_points = [
    "5 seeded notifications (seed data from backend)",
    "Types: ai_prediction_ready · review_reminder",
    "         promotion_recommendation · attendance_alert · system",
    "Priority: high / medium / low",
    "Filter by type and read/unread state",
    "Mark single or all as read",
    "Delete single or clear all",
    "Unread count badge on Sidebar nav icon",
    "Backend: PATCH /notifications/{id}/read",
    "          PATCH /notifications/read-all",
    "          DELETE /notifications/{id}",
    "New predictions auto-create a notification",
]
for i, pt in enumerate(notif_points):
    indent = pt.startswith("         ")
    add_text(s, ("  " if indent else "• ") + pt.strip(),
             0.6, 2.35 + i * 0.41, 5.7, 0.38,
             font_size=10, color=C_SUBTEXT if indent else C_TEXT)

# Settings
add_rect(s, 6.8, 1.72, 6.1, 5.0, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 6.8, 1.72, 6.1, 0.08, fill_rgb=C_PRIMARY)
add_text(s, "⚙️  Settings", 7.0, 1.85, 5.7, 0.4, font_size=14, bold=True, color=C_PRIMARY)

setting_sections = [
    ("General", ["Dark mode toggle (persisted via ThemeContext)",
                 "Compact view toggle",
                 "Language selector (en / fr / es / de)",
                 "Timezone selector"]),
    ("Notifications", ["Email notifications toggle",
                       "Push notifications toggle",
                       "SMS alerts toggle",
                       "Weekly reports toggle"]),
    ("Privacy & Security", ["Placeholder section (future auth hardening)"]),
    ("Account", ["Shows current user role and profile",
                 "Logout button (clears localStorage)"]),
]
y_cursor = 2.35
for section, items in setting_sections:
    add_text(s, section, 7.0, y_cursor, 5.7, 0.32, font_size=11, bold=True, color=C_PRIMARY)
    y_cursor += 0.33
    for item in items:
        add_text(s, "• " + item, 7.0, y_cursor, 5.7, 0.32, font_size=10, color=C_TEXT)
        y_cursor += 0.34
    y_cursor += 0.1

# ─── SLIDE 13 – Deployment ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Deployment")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Deployment — Vercel", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 13 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

# Vercel config card
add_rect(s, 0.4, 1.72, 6.0, 2.5, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 0.4, 1.72, 6.0, 0.08, fill_rgb=C_PRIMARY)
add_text(s, "vercel.json config", 0.6, 1.85, 5.6, 0.4, font_size=14, bold=True, color=C_PRIMARY)
vercel_lines = [
    'buildCommand:   "npm run build"',
    'outputDirectory: "dist"',
    "Python runtime: api/ directory",
    "Serverless function: api/index.py",
]
for i, line in enumerate(vercel_lines):
    add_text(s, line, 0.6, 2.35 + i * 0.46, 5.7, 0.4, font_size=11, color=C_TEXT)

# Env vars card
add_rect(s, 6.8, 1.72, 6.1, 2.5, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 6.8, 1.72, 6.1, 0.08, fill_rgb=C_ACCENT)
add_text(s, "Environment Variables", 7.0, 1.85, 5.7, 0.4, font_size=14, bold=True, color=C_ACCENT)
env_lines = [
    "VITE_API_BASE_URL",
    "  Dev:  http://localhost:5678/api",
    "  Prod: /api  (same-origin Vercel path)",
    "Set in: Vercel Dashboard → Project Settings",
]
for i, line in enumerate(env_lines):
    indent = line.startswith("  ")
    add_text(s, line, 7.0, 2.35 + i * 0.46, 5.7, 0.4,
             font_size=11, color=C_SUBTEXT if indent else C_TEXT)

# Deploy steps
add_rect(s, 0.4, 4.4, 12.5, 2.5, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
add_rect(s, 0.4, 4.4, 12.5, 0.08, fill_rgb=C_EMERALD)
add_text(s, "Deployment Steps", 0.6, 4.53, 12.1, 0.4, font_size=14, bold=True, color=C_EMERALD)
steps = [
    ("1", "npm install -g vercel"),
    ("2", "vercel   (links project)"),
    ("3", "Set VITE_API_BASE_URL in Vercel Dashboard"),
    ("4", "Push to GitHub → auto-deploys on push"),
    ("5", "Confirm /api/health returns { status: 'ok' }"),
]
for i, (num, cmd) in enumerate(steps):
    x = 0.7 + i * 2.45
    add_rect(s, x, 5.05, 0.3, 0.3, fill_rgb=C_EMERALD)
    add_text(s, num, x, 5.05, 0.3, 0.3, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, cmd, x + 0.35, 5.05, 2.0, 0.8, font_size=10, color=C_TEXT)

add_text(s, "⚠  db.json is not durable in serverless deployments. Use a managed database for production.",
         0.6, 6.1, 12.0, 0.35, font_size=10.5, bold=True, color=C_AMBER)

# ─── SLIDE 14 – Known Limitations ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
section_header(s, "Known Limitations")

add_rect(s, 0, 0.55, 13.33, 1.0, fill_rgb=C_PRIMARY)
add_text(s, "Known Limitations", 0.5, 0.6, 10, 0.85, font_size=30, bold=True, color=C_WHITE)
add_text(s, "Slide 14 / 15", 11.3, 0.72, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)

limitations = [
    (C_ROSE,  "Authentication is demo-grade",
     "No password hashing or real token validation. Any password is accepted. "
     "The Bearer token is simply the user's id field stored in localStorage."),
    (C_AMBER, "Mock-first data on most pages",
     "Dashboard, Attendance, Performance, Analytics, Departments, Notifications, and Reports "
     "all render from src/data/mock* files. Backend data is only fetched on Employees and AI Predictions."),
    (C_AMBER, "predictions.ts hardcodes the API URL",
     "src/services/predictions.ts does not read VITE_API_BASE_URL. It always points to "
     "http://localhost:5678/api, which breaks in production without a code change."),
    (C_ROSE,  "JSON persistence is not production-safe",
     "db.json is in-process and file-based. Concurrent requests can cause race conditions. "
     "Vercel serverless instances do not share a file system."),
    (C_AMBER, "Export PDF is a minimal stream",
     "The /reports/export/pdf endpoint generates a basic hand-crafted PDF stream. "
     "It is not a properly formatted report document."),
    (C_SUBTEXT, "No automated tests",
     "Neither frontend (Vitest/Jest) nor backend (pytest) tests are present. "
     "All validation is manual."),
]
for i, (col, title, body) in enumerate(limitations):
    row = i // 2
    col_x = 0.4 + (i % 2) * 6.45
    y = 1.75 + row * 1.8
    add_rect(s, col_x, y, 6.1, 1.65, fill_rgb=C_WHITE, line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.75)
    add_rect(s, col_x, y, 0.06, 1.65, fill_rgb=col)
    add_text(s, title, col_x + 0.2, y + 0.1, 5.8, 0.35, font_size=12, bold=True, color=col)
    add_text(s, body, col_x + 0.2, y + 0.5, 5.8, 1.05, font_size=10, color=C_TEXT)

# ─── SLIDE 15 – Roadmap / Future Improvements ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s, dark=True)

# Header gradient
add_rect(s, 0, 0, 13.33, 1.65, fill_rgb=C_PRIMARY)
add_text(s, "Roadmap & Future Improvements", 0.5, 0.15, 11, 0.9, font_size=32, bold=True, color=C_WHITE)
add_text(s, "Slide 15 / 15", 11.3, 0.25, 1.8, 0.4, font_size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.RIGHT)
add_text(s, "What comes next for AI-HRMS", 0.5, 1.1, 8, 0.45, font_size=14, italic=True,
         color=RGBColor(0xA5, 0xB4, 0xFC))

roadmap = [
    ("🔐", "Secure Authentication",
     "JWT tokens · bcrypt password hashing\nRefresh token rotation · HTTP-only cookies",
     C_PRIMARY),
    ("🗄️", "Production Database",
     "Replace db.json with PostgreSQL or MongoDB\nMigration scripts + ORM (SQLAlchemy / Prisma)",
     C_ACCENT),
    ("🔗", "Full API Integration",
     "Connect Attendance, Performance, Analytics\nand Departments to live backend endpoints",
     C_EMERALD),
    ("🧪", "Automated Testing",
     "Frontend: Vitest + React Testing Library\nBackend: pytest + coverage reporting",
     C_AMBER),
    ("🛡️", "Role-Based Authorization",
     "Route-level permission guards\nAdmin vs HR vs Employee endpoint policies",
     C_ROSE),
    ("📧", "Email & Notifications",
     "Real email delivery for alerts\nWebSocket push for live notification updates",
     RGBColor(0x8B, 0x5C, 0xF6)),
]
cols_r = [roadmap[:3], roadmap[3:]]
col_x_r = [0.35, 6.85]
for ci, col_list in enumerate(cols_r):
    for ri, (icon, title, detail, col) in enumerate(col_list):
        x = col_x_r[ci]
        y = 1.85 + ri * 1.8
        add_rect(s, x, y, 6.1, 1.65, fill_rgb=RGBColor(0x1E, 0x2A, 0x40), line_rgb=col, line_width_pt=1.0)
        add_text(s, icon, x + 0.15, y + 0.18, 0.5, 0.5, font_size=18)
        add_text(s, title, x + 0.72, y + 0.15, 5.2, 0.4, font_size=13, bold=True, color=col)
        add_text(s, detail, x + 0.72, y + 0.6, 5.2, 0.9, font_size=10.5,
                 color=RGBColor(0xCB, 0xD5, 0xE1))

# Footer call to action
add_rect(s, 0, 7.05, 13.33, 0.45, fill_rgb=RGBColor(0x1E, 0x2A, 0x40))
add_text(s, "AI-HRMS  ·  React + TypeScript + Flask + scikit-learn  ·  2026",
         0.4, 7.07, 12.5, 0.35, font_size=10, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)

# ─── Save ─────────────────────────────────────────────────────────────────────
OUT = r"c:\Users\hp\Documents\HR_management_system\HR_Management_System.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
